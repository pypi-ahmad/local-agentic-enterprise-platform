from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import LETTER
from reportlab.pdfgen import canvas

from app.core.config import get_settings


class ReportExportService:
    """Creates report exports in PDF, XLSX, and PPTX formats."""

    def __init__(self) -> None:
        self.settings = get_settings()
        self.output_dir = self.settings.artifacts_dir / "reports"
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def to_pdf(self, filename: str, report_data: dict) -> Path:
        path = self.output_dir / f"{filename}.pdf"
        doc = canvas.Canvas(str(path), pagesize=LETTER)
        doc.setFont("Helvetica", 12)
        y = 740
        doc.drawString(72, y, report_data.get("title", "Business Report"))
        y -= 30
        for key, value in report_data.items():
            doc.drawString(72, y, f"{key}: {value}")
            y -= 20
            if y < 72:
                doc.showPage()
                y = 740
        doc.save()
        return path

    def to_excel(self, filename: str, report_data: dict) -> Path:
        path = self.output_dir / f"{filename}.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "Report"
        ws.append(["Metric", "Value"])
        for key, value in report_data.items():
            ws.append([key, str(value)])
        wb.save(path)
        return path

    def to_pptx(self, filename: str, report_data: dict) -> Path:
        path = self.output_dir / f"{filename}.pptx"
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
        title = slide.shapes.title
        if title:
            title.text = report_data.get("title", "Business Report")
        top = 120
        for key, value in report_data.items():
            textbox = slide.shapes.add_textbox(left=50, top=top, width=620, height=30)
            textbox.text = f"{key}: {value}"
            top += 35
        ppt.save(str(path))
        return path
